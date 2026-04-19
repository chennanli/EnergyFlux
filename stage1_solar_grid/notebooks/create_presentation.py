"""
Generate EnergyFlux Stage 1 presentation (4 slides) as PDF via PPTX.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
import subprocess

OUT_DIR = Path(__file__).parent
DATA_DIR = Path(__file__).parent.parent / "data" / "processed"

# Brand colors
DARK_BLUE = RGBColor(0x04, 0x2C, 0x53)
RED = RGBColor(0xE2, 0x4B, 0x4A)
GREEN = RGBColor(0x1D, 0x9E, 0x75)
BLUE = RGBColor(0x18, 0x5F, 0xA5)
GRAY = RGBColor(0x88, 0x87, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK_BLUE, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=14,
                    color=DARK_BLUE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_header_bar(slide, title):
    """Dark blue header bar across top."""
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    add_text(slide, 0.5, 0.15, 12, 0.6, title, size=28, bold=True,
             color=WHITE)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

add_text(slide, 1, 1.5, 11, 1.2,
         "EnergyFlux", size=48, bold=True, color=WHITE)
add_text(slide, 1, 2.7, 11, 0.8,
         "Physics-Informed Solar Energy Platform with AI-Powered Grid Analytics",
         size=22, color=RGBColor(0xAA, 0xCC, 0xEE))
add_text(slide, 1, 3.8, 11, 0.5,
         "Stage 1: Commercial Behind-the-Meter PV + BESS + Power Flow",
         size=18, color=RGBColor(0x88, 0xAA, 0xCC))

# Key metrics
metrics = [
    "185 kWp Bifacial PV (LONGi LR6-72HBD-385M)",
    "100 kW / 400 kWh BESS with TOU Dispatch",
    "7-Bus 0.4kV Distribution Network — 17,518 Hourly Power Flow Runs",
    "NVIDIA NIM ReAct Agent for Root Cause Analysis",
    "Full REST API (FastAPI + Swagger UI)",
]
add_bullet_list(slide, 1, 4.8, 11, 2.5, metrics, size=16,
                color=RGBColor(0xCC, 0xDD, 0xEE))

add_text(slide, 1, 6.8, 11, 0.4,
         "Chennan Li, PhD, PE  |  github.com/chennanli/EnergyFlux",
         size=14, color=GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: Architecture
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "System Architecture — End-to-End Data Flow")

# Pipeline boxes
pipeline = [
    ("Weather Data\nOpen-Meteo API\n17,518 hours", 0.3),
    ("Irradiance\nForecast\nNHITS 24h", 2.5),
    ("PV Generation\npvlib 4-Mode\n775 MWh/yr", 4.7),
    ("BESS Dispatch\nTOU-Aware\n239 cycles/2yr", 6.9),
    ("Power Flow\npandapower\n17,518 runs", 9.1),
    ("AI Diagnosis\nNVIDIA NIM\nReAct Agent", 11.3),
]
for label, x in pipeline:
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(1.3), Inches(1.8), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE if "Power Flow" in label or "AI" in label else BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.bold = True

# Arrows between boxes
for i in range(5):
    x = 0.3 + (i + 1) * 2.2 - 0.15
    add_text(slide, x - 0.1, 1.5, 0.4, 0.5, "→", size=24, color=BLUE)

# Network diagram (text)
add_text(slide, 0.5, 2.8, 6, 0.5,
         "Network Topology — 0.4kV Commercial Distribution",
         size=16, bold=True)

net_text = """Utility Grid (10kV)  →  Trafo 400kVA (10kV / 0.4kV)  →  LV Busbar
    ├── PV Plant (185 kWp bifacial, LONGi 385W, GCR=0.45)  ·····  300m
    ├── BESS (100 kW / 400 kWh, TOU dispatch)  ·····················  50m
    ├── Biotech Lab (50 kW base 24/7, +15 kW daytime)  ···········  200m
    ├── Office Building (30 kW weekday daytime)  ····················  150m
    └── EV Charging (50 kW weekday daytime — employees)  ·········  250m"""

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.3),
                                  Inches(7), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = net_text
p.font.size = Pt(11)
p.font.color.rgb = DARK_BLUE
p.font.name = "Consolas"

# TOU rates
add_text(slide, 7.8, 2.8, 5, 0.5,
         "BESS TOU Strategy (PG&E California)",
         size=16, bold=True)

tou_items = [
    "Off-Peak (9pm-8am): $0.10/kWh → Charge from grid",
    "Mid-Peak (8am-4pm): $0.20/kWh → Charge from PV surplus",
    "On-Peak (4pm-9pm):  $0.35/kWh → Discharge to avoid import",
    "",
    "Daily cycle: night charge → PV top-up → evening discharge",
    "Result: 239 cycles / 2 years, 96 MWh charged",
]
add_bullet_list(slide, 7.8, 3.3, 5, 3, tou_items, size=12)

# PV 4-mode comparison
add_text(slide, 0.5, 5.8, 6, 0.5,
         "PV 4-Mode Comparison (annual energy)",
         size=14, bold=True)
pv_items = [
    "Mode 1: Fixed tilt 20°           — 663 MWh  (baseline)",
    "Mode 2: Tracker, no backtrack  — 711 MWh  (+7.2%, shading penalty)",
    "Mode 3: Tracker + backtrack     — 751 MWh  (+13.4%)",
    "Mode 4: Bifacial + backtrack    — 775 MWh  (+16.9%, used in power flow)",
]
add_bullet_list(slide, 0.5, 6.2, 7, 1.5, pv_items, size=11, color=DARK_BLUE)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: Power Flow Results
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Power Flow Results — 17,518 Hourly Simulations (2 Years)")

# Main plot
plot_path = DATA_DIR / "powerflow_plot.png"
if plot_path.exists():
    slide.shapes.add_picture(str(plot_path),
                             Inches(0.3), Inches(1.1),
                             Inches(8.5), Inches(5.8))

# Key findings on the right
add_text(slide, 9, 1.2, 4, 0.5, "Key Findings", size=18, bold=True)

findings = [
    "5,363 hours overvoltage (>1.05 pu)",
    "= 31% of all hours, mostly weekends",
    "",
    "Peak voltage: 1.074 pu at PV bus",
    "(ANSI C84.1 limit: 1.05 pu)",
    "",
    "Root cause: PV 135kW peak vs",
    "weekend load 55kW → 80kW surplus",
    "flows back through 300m cable",
    "→ voltage rise at PV bus",
    "",
    "Line loading: max 75% (safe)",
    "No thermal overload risk",
    "",
    "BESS absorbs surplus but",
    "400kWh fills in ~5 hours",
    "→ remaining surplus causes",
    "overvoltage after BESS full",
]
add_bullet_list(slide, 9, 1.7, 4, 5.5, findings, size=11)

# Bottom insight
add_text(slide, 9, 6.3, 4, 0.8,
         "Insight: BESS placement matters more than sizing — "
         "co-locating BESS with PV (DC-coupled) reduces voltage rise "
         "at the source, not just at the busbar.",
         size=11, bold=False, color=RED)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: AI Agent + Next Steps
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "AI-Powered Root Cause Analysis + Roadmap")

# Agent section
add_text(slide, 0.5, 1.2, 6, 0.5,
         "NVIDIA NIM ReAct Agent", size=18, bold=True)

agent_items = [
    "Model: Nemotron-Ultra-253B (via NVIDIA NIM API)",
    "Framework: LangChain ReAct + FAISS RAG",
    "",
    "4 Tools:",
    "  • query_weather() — real-time weather conditions",
    "  • query_pv_status() — PV generation vs expected",
    "  • query_powerflow() — bus voltages, line loading",
    "  • query_knowledge_base() — RAG over engineering docs",
    "",
    "Demo scenario: Midday overvoltage at PV bus",
    "Agent automatically:",
    "  1. Checks weather (clear sky, GHI=930 W/m²)",
    "  2. Checks PV output (135 kW, at capacity)",
    "  3. Checks BESS (SOC=90%, full)",
    "  4. Searches knowledge base",
    "  → Diagnosis: PV surplus + full BESS → reverse flow",
    "  → Actions: increase BESS, add reactive power support",
]
add_bullet_list(slide, 0.5, 1.7, 6, 5.5, agent_items, size=11)

# API section
add_text(slide, 7, 1.2, 6, 0.5,
         "REST API (FastAPI)", size=18, bold=True)

api_items = [
    "GET  /forecast          — 24h GHI forecast",
    "GET  /generation        — PV power (4 modes)",
    "GET  /powerflow          — bus voltages, BESS, loads",
    "GET  /powerflow/summary — violation statistics",
    "POST /diagnose            — trigger RCA agent",
    "",
    "Interactive Swagger UI at /docs",
]
add_bullet_list(slide, 7, 1.7, 5.5, 2.5, api_items, size=12)

# Roadmap
add_text(slide, 7, 4.2, 5.5, 0.5,
         "Roadmap", size=18, bold=True)

roadmap_items = [
    "Stage 2: Chemical plant microgrid (TEP load)",
    "  → 10kV industrial feeder, MW-scale PV + BESS",
    "  → Pyomo LP optimal BESS dispatch",
    "",
    "Stage 3: Data center + 115kV backbone",
    "  → Multiple industrial loads on shared grid",
    "",
    "Stage 4: Utility-scale 100MW solar farm",
    "  → 115kV transmission, grid frequency response",
]
add_bullet_list(slide, 7, 4.7, 5.5, 3, roadmap_items, size=12)

# Tech stack footer
add_text(slide, 0.5, 6.8, 12, 0.4,
         "Tech Stack:  Python · pvlib · pandapower · NeuralForecast · LangChain · NVIDIA NIM · FAISS · FastAPI",
         size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════
pptx_path = OUT_DIR / "EnergyFlux_Stage1_Overview.pptx"
prs.save(str(pptx_path))
print(f"PPTX saved: {pptx_path}")

# Convert to PDF using macOS sips/qlmanage or LibreOffice if available
pdf_path = OUT_DIR / "EnergyFlux_Stage1_Overview.pdf"
try:
    subprocess.run([
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "--headless", "--convert-to", "pdf",
        "--outdir", str(OUT_DIR), str(pptx_path)
    ], check=True, capture_output=True, timeout=30)
    print(f"PDF saved: {pdf_path}")
except Exception:
    print(f"PDF conversion requires LibreOffice. Open the PPTX and export as PDF manually.")
    print(f"Or: brew install --cask libreoffice")
